"""
Media Processor Service.

Handles all non-text media received from WhatsApp:
- Audio: Uazapi transcribes via Whisper → text
- Image: Download base64 → GPT Vision analyzes → description text
- PDF: Download → extract text with PyMuPDF
- XLSX: Download → extract text with openpyxl
- CSV: Download → read as text
- DOCX: Download → extract text with python-docx
- PPTX: Download → extract text with python-pptx

The processed text is then fed to the main conversation agent.
"""

import base64
import csv
import io
import tempfile
from pathlib import Path

import httpx
from loguru import logger
from openai import AsyncOpenAI

from app.core.config import get_settings
from app.services.uazapi import uazapi

settings = get_settings()
openai_client = AsyncOpenAI(api_key=settings.openai_api_key)

# Max characters to extract from documents to avoid token overflow
MAX_DOCUMENT_CHARS = 8000


class MediaProcessor:
    """Process different media types into text for the AI agent."""

    async def process(self, parsed_message: dict) -> str | None:
        """Route media processing based on message type.

        Returns processed text or None if processing fails.
        """
        msg_type = parsed_message["type"]
        message_id = parsed_message["message_id"]

        if not message_id:
            logger.warning("No message_id for media processing")
            return None

        try:
            if msg_type == "audio":
                return await self.process_audio(message_id)
            elif msg_type == "image":
                return await self.process_image(message_id)
            elif msg_type == "document":
                return await self.process_document(parsed_message)
            elif msg_type == "video":
                return "[O usuário enviou um vídeo. Informe que você não consegue processar vídeos no momento.]"
            elif msg_type == "sticker":
                return "[O usuário enviou um sticker/figurinha.]"
            elif msg_type == "location":
                return self._process_location(parsed_message)
            elif msg_type == "contact":
                return "[O usuário compartilhou um contato.]"
            else:
                logger.info(f"Unsupported media type: {msg_type}")
                return None
        except Exception as e:
            logger.error(f"Media processing error ({msg_type}): {e}")
            return None

    # ── Audio ──────────────────────────────────────────────────

    async def process_audio(self, message_id: str) -> str | None:
        """Transcribe audio using Uazapi's built-in Whisper."""
        logger.info(f"Transcribing audio: {message_id}")

        transcription = await uazapi.transcribe_audio(message_id)

        if transcription:
            logger.info(
                f"Audio transcribed ({len(transcription)} chars): "
                f"'{transcription[:80]}...'"
            )
            return f"[Áudio transcrito do usuário]: {transcription}"

        logger.warning(f"Audio transcription failed for {message_id}")
        return "[O usuário enviou um áudio, mas não foi possível transcrevê-lo.]"

    # ── Image ──────────────────────────────────────────────────

    async def process_image(self, message_id: str) -> str | None:
        """Analyze image with GPT Vision."""
        logger.info(f"Processing image: {message_id}")

        media = await uazapi.get_media_base64(message_id)
        if not media:
            return "[O usuário enviou uma imagem, mas não foi possível baixá-la.]"

        base64_data, mimetype = media

        # Determine media type for GPT
        if not mimetype or mimetype == "":
            mimetype = "image/jpeg"

        # Map common mimetypes
        media_type = mimetype.split("/")[-1]
        if media_type not in ("jpeg", "png", "gif", "webp"):
            media_type = "jpeg"

        try:
            response = await openai_client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    "Descreva esta imagem em detalhes, em português. "
                                    "Inclua: o que aparece na imagem, texto visível, "
                                    "cores, objetos, pessoas, cenário. "
                                    "Se houver texto/documento na imagem, transcreva-o."
                                ),
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/{media_type};base64,{base64_data}",
                                    "detail": "high",
                                },
                            },
                        ],
                    }
                ],
                max_tokens=1024,
                temperature=0.2,
            )

            description = response.choices[0].message.content or ""
            logger.info(f"Image analyzed ({len(description)} chars)")
            return f"[O usuário enviou uma imagem. Descrição da imagem]: {description}"

        except Exception as e:
            logger.error(f"GPT Vision error: {e}")
            return "[O usuário enviou uma imagem, mas ocorreu um erro ao analisá-la.]"

    # ── Documents ──────────────────────────────────────────────

    async def process_document(self, parsed_message: dict) -> str | None:
        """Process document based on mimetype/extension."""
        message_id = parsed_message["message_id"]
        content = parsed_message.get("content", {})

        # Try to determine file type from content metadata
        mimetype = ""
        filename = ""

        if isinstance(content, dict):
            doc_msg = (
                content.get("documentMessage", {})
                or content.get("documentWithCaptionMessage", {}).get(
                    "message", {}
                ).get("documentMessage", {})
            )
            mimetype = doc_msg.get("mimetype", "")
            filename = doc_msg.get("fileName", "")

        logger.info(f"Processing document: {filename} ({mimetype})")

        # Determine processor by mimetype or filename extension
        ext = Path(filename).suffix.lower() if filename else ""

        if mimetype == "application/pdf" or ext == ".pdf":
            return await self._process_pdf(message_id, filename)
        elif mimetype in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or ext in (".xlsx", ".xls"):
            return await self._process_xlsx(message_id, filename)
        elif ext == ".csv" or mimetype == "text/csv":
            return await self._process_csv(message_id, filename)
        elif mimetype in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or ext in (".docx", ".doc"):
            return await self._process_docx(message_id, filename)
        elif mimetype in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ) or ext in (".pptx",):
            return await self._process_pptx(message_id, filename)
        elif mimetype.startswith("text/") or ext in (".txt", ".json", ".xml", ".md"):
            return await self._process_text_file(message_id, filename)
        else:
            return (
                f"[O usuário enviou um documento: {filename or 'arquivo'}. "
                f"Tipo: {mimetype or 'desconhecido'}. "
                f"Não é possível processar este formato no momento.]"
            )

    async def _download_file_bytes(self, message_id: str) -> bytes | None:
        """Download file as bytes via base64."""
        media = await uazapi.get_media_base64(message_id)
        if not media:
            return None
        base64_data, _ = media
        return base64.b64decode(base64_data)

    async def _process_pdf(self, message_id: str, filename: str) -> str:
        """Extract text from PDF using PyMuPDF."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou o PDF '{filename}', mas não foi possível baixá-lo.]"

        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text_parts = []
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(f"--- Página {page_num} ---\n{page_text.strip()}")
            doc.close()

            full_text = "\n\n".join(text_parts)

            if not full_text.strip():
                # PDF might be scanned/image-based — try GPT Vision
                return await self._process_pdf_as_image(message_id, filename)

            if len(full_text) > MAX_DOCUMENT_CHARS:
                full_text = (
                    full_text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... documento truncado, total de {len(full_text)} caracteres]"
                )

            return (
                f"[O usuário enviou o PDF '{filename}'. "
                f"Conteúdo extraído]:\n\n{full_text}"
            )

        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return f"[O usuário enviou o PDF '{filename}', mas houve erro ao extrair o conteúdo.]"

    async def _process_pdf_as_image(self, message_id: str, filename: str) -> str:
        """For scanned PDFs, send first page as image to GPT Vision."""
        try:
            import fitz

            file_bytes = await self._download_file_bytes(message_id)
            if not file_bytes:
                return f"[PDF '{filename}' sem texto extraível e não foi possível baixar.]"

            doc = fitz.open(stream=file_bytes, filetype="pdf")
            if len(doc) == 0:
                return f"[PDF '{filename}' está vazio.]"

            # Render first page as image
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            img_bytes = pix.tobytes("png")
            doc.close()

            b64 = base64.b64encode(img_bytes).decode()

            response = await openai_client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    "Este é um PDF escaneado. Transcreva todo o texto "
                                    "visível nesta página, em português. Mantenha a "
                                    "estrutura e formatação o mais próximo possível."
                                ),
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{b64}",
                                    "detail": "high",
                                },
                            },
                        ],
                    }
                ],
                max_tokens=2048,
                temperature=0.1,
            )

            ocr_text = response.choices[0].message.content or ""
            return (
                f"[O usuário enviou o PDF '{filename}' (escaneado). "
                f"Texto extraído via OCR]:\n\n{ocr_text}"
            )

        except Exception as e:
            logger.error(f"PDF-as-image error: {e}")
            return f"[PDF '{filename}' parece ser escaneado e não foi possível processar.]"

    async def _process_xlsx(self, message_id: str, filename: str) -> str:
        """Extract data from XLSX."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou '{filename}', mas não foi possível baixá-lo.]"

        try:
            from openpyxl import load_workbook

            wb = load_workbook(filename=io.BytesIO(file_bytes), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cell_values = [str(c) if c is not None else "" for c in row]
                    if any(v.strip() for v in cell_values):
                        rows.append(" | ".join(cell_values))

                if rows:
                    text_parts.append(
                        f"--- Aba: {sheet_name} ---\n" + "\n".join(rows)
                    )

            wb.close()
            full_text = "\n\n".join(text_parts)

            if len(full_text) > MAX_DOCUMENT_CHARS:
                full_text = (
                    full_text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... planilha truncada, total de {len(full_text)} caracteres]"
                )

            return (
                f"[O usuário enviou a planilha '{filename}'. "
                f"Conteúdo extraído]:\n\n{full_text}"
            )

        except Exception as e:
            logger.error(f"XLSX extraction error: {e}")
            return f"[O usuário enviou '{filename}', mas houve erro ao processar.]"

    async def _process_csv(self, message_id: str, filename: str) -> str:
        """Read CSV content."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou '{filename}', mas não foi possível baixá-lo.]"

        try:
            # Try utf-8 then latin-1
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                text = file_bytes.decode("latin-1")

            reader = csv.reader(io.StringIO(text))
            rows = [" | ".join(row) for row in reader]
            full_text = "\n".join(rows)

            if len(full_text) > MAX_DOCUMENT_CHARS:
                full_text = (
                    full_text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... CSV truncado, total de {len(full_text)} caracteres]"
                )

            return (
                f"[O usuário enviou o CSV '{filename}'. "
                f"Conteúdo]:\n\n{full_text}"
            )

        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            return f"[O usuário enviou '{filename}', mas houve erro ao processar.]"

    async def _process_docx(self, message_id: str, filename: str) -> str:
        """Extract text from DOCX."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou '{filename}', mas não foi possível baixá-lo.]"

        try:
            from docx import Document

            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)

            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        full_text += "\n" + " | ".join(cells)

            if len(full_text) > MAX_DOCUMENT_CHARS:
                full_text = (
                    full_text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... documento truncado, total de {len(full_text)} caracteres]"
                )

            return (
                f"[O usuário enviou o documento '{filename}'. "
                f"Conteúdo extraído]:\n\n{full_text}"
            )

        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return f"[O usuário enviou '{filename}', mas houve erro ao processar.]"

    async def _process_pptx(self, message_id: str, filename: str) -> str:
        """Extract text from PPTX."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou '{filename}', mas não foi possível baixá-lo.]"

        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_parts.append(
                        f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
                    )

            full_text = "\n\n".join(text_parts)

            if len(full_text) > MAX_DOCUMENT_CHARS:
                full_text = (
                    full_text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... apresentação truncada, total de {len(full_text)} caracteres]"
                )

            return (
                f"[O usuário enviou a apresentação '{filename}'. "
                f"Conteúdo extraído]:\n\n{full_text}"
            )

        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return f"[O usuário enviou '{filename}', mas houve erro ao processar.]"

    async def _process_text_file(self, message_id: str, filename: str) -> str:
        """Read plain text files (txt, json, xml, md)."""
        file_bytes = await self._download_file_bytes(message_id)
        if not file_bytes:
            return f"[O usuário enviou '{filename}', mas não foi possível baixá-lo.]"

        try:
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                text = file_bytes.decode("latin-1")

            if len(text) > MAX_DOCUMENT_CHARS:
                text = (
                    text[:MAX_DOCUMENT_CHARS]
                    + f"\n\n[... arquivo truncado, total de {len(text)} caracteres]"
                )

            return (
                f"[O usuário enviou o arquivo '{filename}'. "
                f"Conteúdo]:\n\n{text}"
            )

        except Exception as e:
            logger.error(f"Text file error: {e}")
            return f"[O usuário enviou '{filename}', mas houve erro ao processar.]"

    def _process_location(self, parsed_message: dict) -> str:
        """Extract location data."""
        content = parsed_message.get("content", {})
        loc = content.get("locationMessage", {})
        lat = loc.get("degreesLatitude", "")
        lng = loc.get("degreesLongitude", "")
        name = loc.get("name", "")
        address = loc.get("address", "")

        parts = [f"[O usuário compartilhou uma localização]"]
        if name:
            parts.append(f"Local: {name}")
        if address:
            parts.append(f"Endereço: {address}")
        if lat and lng:
            parts.append(f"Coordenadas: {lat}, {lng}")

        return " | ".join(parts)


# Singleton
media_processor = MediaProcessor()
